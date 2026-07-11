"""
AI Command Center - External Tools MCP Server
------------------------------------------------
هدف هذا السيرفر: تنفيذ العمليات التي لا يقدر نموذج اللغة ينفذها بنفسه
(فتح روابط فيديو فعلياً، فك ضغط أرشيفات) ويرجّع بيانات حقيقية فقط.
مصمم للنشر المجاني على Render.com (Web Service) ثم ربطه كأداة MCP خارجية
داخل إعدادات تطبيق genspark (Tools / MCP Servers / Integrations).

التشغيل محلياً للاختبار:
    pip install -r requirements.txt
    python server.py
ثم اختبر بمتصفح: GET /health

النشر على Render:
    Build Command:  pip install -r requirements.txt
    Start Command:  python server.py
    (Render يمرر متغير البيئة PORT تلقائياً، السيرفر يقرأه بنفسه)
"""

import os
import re
import shutil
import tempfile
import zipfile
from pathlib import Path

import requests
from mcp.server.fastmcp import FastMCP

mcp = FastMCP(
    "ai-command-center-tools",
    host="0.0.0.0",
    port=int(os.environ.get("PORT", 8000)),
)

MAX_FILE_MB = 100
YOUTUBE_TIKTOK_INSTAGRAM = re.compile(
    r"(youtu\.be/|youtube\.com/(watch|shorts)|tiktok\.com/|instagram\.com/reel)"
)


def _download_to_temp(url: str, tmp_dir: str, max_mb: int = MAX_FILE_MB) -> tuple[str | None, dict | None]:
    """
    يحمّل ملف من رابط لمجلد مؤقت. يرجع (المسار, None) عند النجاح
    أو (None, {"success": False, "error": ...}) عند الفشل.
    """
    local_path = os.path.join(tmp_dir, "downloaded_file")
    try:
        r = requests.get(url, timeout=60, stream=True)
        r.raise_for_status()
        size = 0
        with open(local_path, "wb") as f:
            for chunk in r.iter_content(chunk_size=1024 * 256):
                size += len(chunk)
                if size > max_mb * 1024 * 1024:
                    return None, {"success": False, "error": f"الملف أكبر من الحد الأقصى المسموح ({max_mb}MB)."}
                f.write(chunk)
    except Exception as e:
        return None, {"success": False, "error": f"تعذّر تحميل الملف: {e}"}
    return local_path, None


# ----------------------------------------------------------------------
# أداة 1: تحليل رابط فيديو (يوتيوب / تيك توك / إنستغرام)
# ----------------------------------------------------------------------
def _analyze_video_url_impl(url: str) -> dict:
    if not YOUTUBE_TIKTOK_INSTAGRAM.search(url):
        return {"success": False, "error": "الرابط غير مدعوم حالياً (يوتيوب/تيك توك/إنستغرام فقط)."}

    try:
        import yt_dlp
    except ImportError:
        return {"success": False, "error": "مكتبة yt-dlp غير مثبتة على السيرفر."}

    ydl_opts = {
        "quiet": True,
        "skip_download": True,
        "writesubtitles": True,
        "writeautomaticsub": True,
        "subtitleslangs": ["ar", "en"],
    }
    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)
    except Exception as e:
        return {"success": False, "error": f"تعذّر جلب الفيديو: {e}"}

    transcript_text = _extract_transcript(info)

    return {
        "success": True,
        "title": info.get("title"),
        "uploader": info.get("uploader"),
        "duration_seconds": info.get("duration"),
        "description": (info.get("description") or "")[:2000],
        "transcript_available": bool(transcript_text),
        "transcript": transcript_text[:8000] if transcript_text else None,
    }


@mcp.tool()
def analyze_video_url(url: str) -> dict:
    """
    يجلب بيانات حقيقية عن رابط فيديو: العنوان، الوصف، القناة، المدة،
    والترانسكريبت إن توفر (captions أو auto-captions).
    لا يخمّن أبداً: إذا فشل الجلب يرجع success=False مع سبب واضح.
    """
    return _analyze_video_url_impl(url)


def _extract_transcript(info: dict) -> str | None:
    """يحاول استخراج نص الترجمة من subtitles أو automatic_captions إن وُجدت."""
    for key in ("subtitles", "automatic_captions"):
        tracks = info.get(key) or {}
        for lang in ("ar", "en"):
            if lang in tracks:
                for fmt in tracks[lang]:
                    if fmt.get("ext") in ("vtt", "srv1", "srt") and fmt.get("url"):
                        try:
                            r = requests.get(fmt["url"], timeout=15)
                            if r.ok:
                                return _clean_subtitle_text(r.text)
                        except Exception:
                            continue
    return None


def _clean_subtitle_text(raw: str) -> str:
    lines = []
    for line in raw.splitlines():
        line = line.strip()
        if not line or "-->" in line or line.startswith(("WEBVTT", "Kind:", "Language:")):
            continue
        if line.isdigit():
            continue
        lines.append(re.sub(r"<[^>]+>", "", line))
    seen, out = set(), []
    for l in lines:
        if l not in seen:
            seen.add(l)
            out.append(l)
    return " ".join(out)


# ----------------------------------------------------------------------
# أداة 2: فك ضغط أرشيف (ZIP / RAR / 7z) من رابط وإرجاع محتواه الفعلي
# ----------------------------------------------------------------------
def _extract_archive_impl(file_url: str) -> dict:
    tmp_dir = tempfile.mkdtemp()
    try:
        archive_path, err = _download_to_temp(file_url, tmp_dir, max_mb=MAX_FILE_MB)
        if err:
            return err

        extract_dir = os.path.join(tmp_dir, "extracted")
        os.makedirs(extract_dir, exist_ok=True)
        infos = {}
        archive_type = None

        try:
            if zipfile.is_zipfile(archive_path):
                archive_type = "zip"
                with zipfile.ZipFile(archive_path) as zf:
                    infos = {i.filename: i.file_size for i in zf.infolist() if not i.is_dir()}
                    zf.extractall(extract_dir)
            else:
                import rarfile
                import py7zr
                if rarfile.is_rarfile(archive_path):
                    archive_type = "rar"
                    with rarfile.RarFile(archive_path) as rf:
                        infos = {i.filename: i.file_size for i in rf.infolist() if not i.is_dir()}
                        rf.extractall(extract_dir)
                elif py7zr.is_7zfile(archive_path):
                    archive_type = "7z"
                    with py7zr.SevenZipFile(archive_path, mode="r") as zf:
                        zf.extractall(extract_dir)
                    for root, _, files in os.walk(extract_dir):
                        for fn in files:
                            full = os.path.join(root, fn)
                            infos[os.path.relpath(full, extract_dir)] = os.path.getsize(full)
                else:
                    return {"success": False, "error": "صيغة الأرشيف غير مدعومة (المدعوم حالياً: ZIP, RAR, 7z)."}
        except Exception as e:
            return {"success": False, "error": f"تعذّر فتح/فك الأرشيف: {e}"}

        results = []
        for name, size in infos.items():
            size_mb = size / (1024 * 1024)
            entry = {"name": name, "size_mb": round(size_mb, 2)}
            if size_mb > MAX_FILE_MB:
                entry["note"] = "تم تجاوز الحد الأقصى 100MB، لم تتم معالجته."
                results.append(entry)
                continue
            try:
                full_path = Path(extract_dir) / name
                if full_path.exists() and full_path.suffix.lower() in (".txt", ".md", ".json", ".csv", ".py", ".js", ".html"):
                    entry["excerpt"] = full_path.read_text(errors="ignore")[:500]
            except Exception as e:
                entry["note"] = f"فشل قراءة المحتوى: {e}"
            results.append(entry)

        return {"success": True, "archive_type": archive_type, "file_count": len(results), "files": results}
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


@mcp.tool()
def extract_archive(file_url: str) -> dict:
    """
    يحمّل أرشيف (ZIP أو RAR أو 7z) من رابط ويفك ضغطه فعلياً، ويرجع قائمة
    الملفات الداخلية مع مقتطف نصي من كل ملف نصي (حد أقصى 100MB لكل ملف).
    """
    return _extract_archive_impl(file_url)


# ----------------------------------------------------------------------
# أداة 3: تحليل صورة (أبعاد، صيغة، بيانات EXIF)
# ----------------------------------------------------------------------
def _analyze_image_impl(url: str) -> dict:
    tmp_dir = tempfile.mkdtemp()
    try:
        path, err = _download_to_temp(url, tmp_dir, max_mb=30)
        if err:
            return err
        try:
            from PIL import Image, ExifTags
            with Image.open(path) as img:
                exif_data = {}
                raw_exif = img.getexif()
                if raw_exif:
                    for tag_id, value in raw_exif.items():
                        tag = ExifTags.TAGS.get(tag_id, str(tag_id))
                        if isinstance(value, (str, int, float)):
                            exif_data[tag] = value
                return {
                    "success": True,
                    "format": img.format,
                    "width": img.width,
                    "height": img.height,
                    "mode": img.mode,
                    "file_size_kb": round(os.path.getsize(path) / 1024, 1),
                    "exif": exif_data or None,
                }
        except Exception as e:
            return {"success": False, "error": f"تعذّر قراءة الصورة (صيغة غير مدعومة أو ملف تالف): {e}"}
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


@mcp.tool()
def analyze_image_url(url: str) -> dict:
    """
    يجلب صورة من رابط ويرجع بياناتها الحقيقية: الأبعاد، الصيغة، حجم الملف،
    وبيانات EXIF إن وُجدت (مثل نوع الكاميرا وتاريخ الالتقاط).
    لا يصف محتوى الصورة بصرياً — هذا يعتمد على قدرة النموذج البصرية نفسه.
    """
    return _analyze_image_impl(url)


def _ocr_image_impl(url: str) -> dict:
    tmp_dir = tempfile.mkdtemp()
    try:
        path, err = _download_to_temp(url, tmp_dir, max_mb=30)
        if err:
            return err
        try:
            import pytesseract
            from PIL import Image
            text = pytesseract.image_to_string(Image.open(path), lang="ara+eng").strip()
            return {"success": True, "text": text[:8000], "has_text": bool(text)}
        except Exception as e:
            return {"success": False, "error": f"تعذّر تشغيل OCR: {e}"}
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


@mcp.tool()
def ocr_image_url(url: str) -> dict:
    """
    يستخرج فعلياً أي نص مكتوب داخل صورة (عربي أو إنجليزي) عبر Tesseract OCR.
    مفيد للقطات شاشة أو مستندات ممسوحة ضوئياً أو صور فيها نص.
    لا يخمّن النص — إذا ما وجد نص يرجع has_text: false.
    """
    return _ocr_image_impl(url)


# ----------------------------------------------------------------------
# أداة 4: تحليل ملف صوتي (مدة، جودة، بيانات وصفية)
# ----------------------------------------------------------------------
def _analyze_audio_impl(url: str) -> dict:
    tmp_dir = tempfile.mkdtemp()
    try:
        path, err = _download_to_temp(url, tmp_dir, max_mb=60)
        if err:
            return err
        try:
            import mutagen
            f = mutagen.File(path)
            if f is None:
                return {"success": False, "error": "الصيغة غير مدعومة أو الملف ليس ملف صوتي صالح."}
            tags = {}
            if f.tags:
                for key in ("title", "artist", "album", "date", "TIT2", "TPE1", "TALB"):
                    if key in f.tags:
                        tags[key] = str(f.tags[key])
            return {
                "success": True,
                "duration_seconds": round(f.info.length, 1) if hasattr(f.info, "length") else None,
                "bitrate_kbps": round(f.info.bitrate / 1000, 1) if hasattr(f.info, "bitrate") else None,
                "sample_rate": getattr(f.info, "sample_rate", None),
                "file_size_kb": round(os.path.getsize(path) / 1024, 1),
                "tags": tags or None,
                "note": "هذه بيانات وصفية فقط. تفريغ الصوت لنص (Speech-to-Text) غير مفعّل حالياً لأنه يحتاج موارد أكبر من الخطة المجانية.",
            }
        except Exception as e:
            return {"success": False, "error": f"تعذّر قراءة الملف الصوتي: {e}"}
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


@mcp.tool()
def analyze_audio_url(url: str) -> dict:
    """
    يجلب ملف صوتي من رابط ويرجع بياناته الحقيقية: المدة، معدل البت،
    والبيانات الوصفية (فنان/عنوان إن وُجدت). لا يفرّغ الكلام إلى نص.
    """
    return _analyze_audio_impl(url)


# ----------------------------------------------------------------------
# أداة 5: قراءة نص من مستند (PDF / DOCX / TXT / MD / CSV / JSON)
# ----------------------------------------------------------------------
def _read_document_impl(url: str) -> dict:
    tmp_dir = tempfile.mkdtemp()
    try:
        path, err = _download_to_temp(url, tmp_dir, max_mb=50)
        if err:
            return err

        ext = Path(url.split("?")[0]).suffix.lower()
        try:
            if ext == ".pdf":
                from pypdf import PdfReader
                reader = PdfReader(path)
                text = "\n".join((p.extract_text() or "") for p in reader.pages[:30])
                return {
                    "success": True,
                    "type": "pdf",
                    "page_count": len(reader.pages),
                    "text": text[:12000],
                    "truncated": len(text) > 12000,
                }
            elif ext == ".docx":
                import docx
                d = docx.Document(path)
                text = "\n".join(p.text for p in d.paragraphs if p.text.strip())
                return {
                    "success": True,
                    "type": "docx",
                    "paragraph_count": len(d.paragraphs),
                    "text": text[:12000],
                    "truncated": len(text) > 12000,
                }
            elif ext in (".txt", ".md", ".csv", ".json", ".py", ".js", ".html", ".xml"):
                text = Path(path).read_text(errors="ignore")
                return {"success": True, "type": ext.lstrip("."), "text": text[:12000], "truncated": len(text) > 12000}
            elif ext == ".xlsx":
                import io
                import openpyxl
                with open(path, "rb") as fh:
                    buf = io.BytesIO(fh.read())
                wb = openpyxl.load_workbook(buf, data_only=True, read_only=True)
                sheets_preview = {}
                for sheet_name in wb.sheetnames[:5]:
                    ws = wb[sheet_name]
                    rows = []
                    for i, row in enumerate(ws.iter_rows(values_only=True)):
                        if i >= 50:
                            break
                        rows.append([("" if c is None else str(c)) for c in row])
                    sheets_preview[sheet_name] = rows
                return {
                    "success": True,
                    "type": "xlsx",
                    "sheet_names": wb.sheetnames,
                    "sheets_preview": sheets_preview,
                    "note": "معاينة أول 50 صف من أول 5 شيتات فقط.",
                }
            elif ext == ".pptx":
                import pptx
                prs = pptx.Presentation(path)
                slides = []
                for i, slide in enumerate(prs.slides):
                    texts = [
                        shape.text_frame.text
                        for shape in slide.shapes
                        if shape.has_text_frame and shape.text_frame.text.strip()
                    ]
                    slides.append({"slide": i + 1, "text": "\n".join(texts)})
                return {"success": True, "type": "pptx", "slide_count": len(slides), "slides": slides}
            else:
                return {"success": False, "error": f"صيغة الملف '{ext or 'غير معروفة'}' غير مدعومة حالياً (مدعوم: pdf, docx, xlsx, pptx, txt, md, csv, json)."}
        except Exception as e:
            return {"success": False, "error": f"تعذّر قراءة المستند: {e}"}
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


@mcp.tool()
def read_document_url(url: str) -> dict:
    """
    يحمّل مستند من رابط (PDF, DOCX, TXT, MD, CSV, JSON) ويستخرج نصه الحقيقي.
    الرابط يجب أن ينتهي بامتداد الملف الصحيح ليعرف نوعه.
    """
    return _read_document_impl(url)


# ----------------------------------------------------------------------
# أداة 6: تحليل ملف APK (اسم الحزمة، الصلاحيات، الإصدار)
# ----------------------------------------------------------------------
def _analyze_apk_impl(url: str) -> dict:
    tmp_dir = tempfile.mkdtemp()
    try:
        path, err = _download_to_temp(url, tmp_dir, max_mb=100)
        if err:
            return err
        try:
            from androguard.core.apk import APK
            apk = APK(path)
            return {
                "success": True,
                "package_name": apk.get_package(),
                "app_name": apk.get_app_name(),
                "version_name": apk.get_androidversion_name(),
                "version_code": apk.get_androidversion_code(),
                "min_sdk": apk.get_min_sdk_version(),
                "target_sdk": apk.get_target_sdk_version(),
                "permissions": apk.get_permissions()[:50],
                "permission_count": len(apk.get_permissions()),
                "is_signed": apk.is_signed(),
            }
        except Exception as e:
            return {"success": False, "error": f"تعذّر تحليل ملف APK (قد يكون تالفاً أو الرابط لا يشير لملف APK حقيقي): {e}"}
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


@mcp.tool()
def analyze_apk_url(url: str) -> dict:
    """
    يحمّل ملف APK من رابط ويحلله فعلياً: اسم الحزمة، اسم التطبيق، الإصدار،
    قائمة الصلاحيات الحقيقية، الحد الأدنى/الأقصى لإصدار أندرويد.
    """
    return _analyze_apk_impl(url)


# ----------------------------------------------------------------------
# نقاط REST عادية (GET) — لتُستخدم إذا كان لدى genspark ميزة عامة
# "قراءة رابط" تفتح أي URL وتقرأ محتواه، بدل بروتوكول MCP الكامل.
# مثال: GET /api/analyze-video?url=https://youtu.be/xxxx
# ----------------------------------------------------------------------
@mcp.custom_route("/api/analyze-video", methods=["GET"])
async def api_analyze_video(request):
    from starlette.responses import JSONResponse
    url = request.query_params.get("url")
    if not url:
        return JSONResponse({"success": False, "error": "مطلوب باراميتر url"}, status_code=400)
    return JSONResponse(_analyze_video_url_impl(url))


@mcp.custom_route("/api/extract-archive", methods=["GET"])
async def api_extract_archive(request):
    from starlette.responses import JSONResponse
    url = request.query_params.get("url")
    if not url:
        return JSONResponse({"success": False, "error": "مطلوب باراميتر url"}, status_code=400)
    return JSONResponse(_extract_archive_impl(url))


@mcp.custom_route("/api/analyze-image", methods=["GET"])
async def api_analyze_image(request):
    from starlette.responses import JSONResponse
    url = request.query_params.get("url")
    if not url:
        return JSONResponse({"success": False, "error": "مطلوب باراميتر url"}, status_code=400)
    return JSONResponse(_analyze_image_impl(url))


@mcp.custom_route("/api/ocr-image", methods=["GET"])
async def api_ocr_image(request):
    from starlette.responses import JSONResponse
    url = request.query_params.get("url")
    if not url:
        return JSONResponse({"success": False, "error": "مطلوب باراميتر url"}, status_code=400)
    return JSONResponse(_ocr_image_impl(url))


@mcp.custom_route("/api/analyze-audio", methods=["GET"])
async def api_analyze_audio(request):
    from starlette.responses import JSONResponse
    url = request.query_params.get("url")
    if not url:
        return JSONResponse({"success": False, "error": "مطلوب باراميتر url"}, status_code=400)
    return JSONResponse(_analyze_audio_impl(url))


@mcp.custom_route("/api/read-document", methods=["GET"])
async def api_read_document(request):
    from starlette.responses import JSONResponse
    url = request.query_params.get("url")
    if not url:
        return JSONResponse({"success": False, "error": "مطلوب باراميتر url"}, status_code=400)
    return JSONResponse(_read_document_impl(url))


@mcp.custom_route("/api/analyze-apk", methods=["GET"])
async def api_analyze_apk(request):
    from starlette.responses import JSONResponse
    url = request.query_params.get("url")
    if not url:
        return JSONResponse({"success": False, "error": "مطلوب باراميتر url"}, status_code=400)
    return JSONResponse(_analyze_apk_impl(url))


# ----------------------------------------------------------------------
# نقطة فحص صحة السيرفر (يُستخدم للتأكد أنه يعمل قبل ربطه بـ genspark)
# ----------------------------------------------------------------------
@mcp.custom_route("/health", methods=["GET"])
async def health(request):
    from starlette.responses import JSONResponse
    return JSONResponse({
        "status": "ok",
        "mcp_tools": [
            "analyze_video_url", "extract_archive", "analyze_image_url", "ocr_image_url",
            "analyze_audio_url", "read_document_url", "analyze_apk_url",
        ],
        "rest_endpoints": [
            "/api/analyze-video?url=...", "/api/extract-archive?url=...",
            "/api/analyze-image?url=...", "/api/ocr-image?url=...",
            "/api/analyze-audio?url=...", "/api/read-document?url=...",
            "/api/analyze-apk?url=...",
        ],
    })


if __name__ == "__main__":
    # streamable-http = بروتوكول MCP عبر HTTP، هذا ما تحتاجه genspark كسيرفر MCP خارجي (remote)
    # ملاحظة: host/port يُقرآن من إعداد FastMCP نفسه أعلاه (وليس من run())
    mcp.run(transport="streamable-http")
